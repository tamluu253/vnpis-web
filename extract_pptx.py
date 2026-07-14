import json
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def get_col(x):
    if x < 1000000: return 0
    if x < 3000000: return 1
    if x < 5000000: return 2
    if x < 6500000: return 3
    return 4

def get_row(y):
    if y < 2000000: return 0
    if y < 4500000: return 1
    if y < 7000000: return 2
    return 3

def extract_and_match():
    data_path = 'src/data/sj-printers.json'
    pptx_path = 'public/images/pad-printers/sj/Catalog and price.pptx'
    out_dir = 'public/images/pad-printers/sj'
    
    with open(data_path, 'r', encoding='utf-8') as f:
        printers = json.load(f)
        
    prs = Presentation(pptx_path)
    
    def find_pictures(shapes):
        pics = []
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pics.append(shape)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                pics.extend(find_pictures(shape.shapes))
        return pics

    # Process Slide 3 (Models 1-20, which is index 0-19)
    slide3_pics = find_pictures(prs.slides[2].shapes)
    for pic in slide3_pics:
        row = get_row(pic.top)
        col = get_col(pic.left)
        model_idx = row * 5 + col
        if model_idx < len(printers):
            model_name = printers[model_idx]['model']
            
            image_bytes = pic.image.blob
            ext = pic.image.ext
            filename = f"{model_name.lower().replace('-', '_')}.{ext}"
            filepath = os.path.join(out_dir, filename)
            
            with open(filepath, 'wb') as f:
                f.write(image_bytes)
                
            printers[model_idx]['image'] = f"/images/pad-printers/sj/{filename}"

    # Process Slide 4 (Models 21-30, which is index 20-29)
    slide4_pics = find_pictures(prs.slides[3].shapes)
    for pic in slide4_pics:
        row = get_row(pic.top)
        col = get_col(pic.left)
        model_idx = 20 + row * 5 + col
        if model_idx < len(printers):
            model_name = printers[model_idx]['model']
            
            image_bytes = pic.image.blob
            ext = pic.image.ext
            filename = f"{model_name.lower().replace('-', '_')}.{ext}"
            filepath = os.path.join(out_dir, filename)
            
            with open(filepath, 'wb') as f:
                f.write(image_bytes)
                
            printers[model_idx]['image'] = f"/images/pad-printers/sj/{filename}"
            
    with open(data_path, 'w', encoding='utf-8') as f:
        json.dump(printers, f, ensure_ascii=False, indent=2)
        
    print("Done! Extracted images by grid layout.")

if __name__ == '__main__':
    extract_and_match()
